"""Converter for PPTX files using python-pptx."""

from typing import Optional

from pptx import Presentation

from office2md.converters.base_converter import BaseConverter


class PptxConverter(BaseConverter):
    """Converter for PPTX files to Markdown."""

    def __init__(
        self,
        input_path: str,
        output_path: Optional[str] = None,
        include_notes: bool = True,
    ):
        """
        Initialize the PPTX converter.

        Args:
            input_path: Path to the input PPTX file
            output_path: Optional path for the output Markdown file
            include_notes: If True, include speaker notes in the output
        """
        super().__init__(input_path, output_path)
        self.include_notes = include_notes

    def convert(self) -> str:
        """
        Convert PPTX to Markdown.

        Returns:
            The Markdown content as a string
        """
        self.logger.info(f"Converting PPTX file: {self.input_path}")
        prs = Presentation(self.input_path)
        markdown_lines = []

        for slide_num, slide in enumerate(prs.slides, 1):
            markdown_lines.append(f"## Slide {slide_num}")
            markdown_lines.append("")

            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    # Check if it's a title (usually first text element)
                    if shape == slide.shapes[0] and hasattr(shape, "text_frame"):
                        markdown_lines.append(f"### {text}")
                    else:
                        # Split by lines and format as list or paragraph
                        lines = text.split("\n")
                        if len(lines) > 1:
                            for line in lines:
                                line = line.strip()
                                if line:
                                    markdown_lines.append(f"- {line}")
                        else:
                            markdown_lines.append(text)
                    markdown_lines.append("")

                # Handle tables
                if shape.has_table:
                    table = shape.table
                    for i, row in enumerate(table.rows):
                        cells = [cell.text.strip() for cell in row.cells]
                        markdown_lines.append("| " + " | ".join(cells) + " |")
                        if i == 0:  # Add header separator
                            markdown_lines.append(
                                "| " + " | ".join(["---"] * len(cells)) + " |"
                            )
                    markdown_lines.append("")

            # Add speaker notes if available
            if self.include_notes and slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    markdown_lines.append("**Notes:**")
                    markdown_lines.append("")
                    markdown_lines.append(notes_text)
                    markdown_lines.append("")

            markdown_lines.append("---")
            markdown_lines.append("")

        return "\n".join(markdown_lines)
